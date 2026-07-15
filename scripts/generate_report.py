import sys
import os
import json
import base64
import io
import copy
from pptx import Presentation

def clone_slide(prs, src_slide):
    # Create new slide using the same layout
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    # Copy all shapes
    for shape in src_slide.shapes:
        new_shape_xml = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_shape_xml)
    return new_slide

def main():
    if len(sys.argv) < 3:
        print("Usage: generate_report.py <input_json_path> <output_pptx_path>")
        sys.exit(1)

    json_path = sys.argv[1]
    output_path = sys.argv[2]

    if not os.path.exists(json_path):
        print(f"Error: Input JSON file not found: {json_path}")
        sys.exit(1)

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    report_date = data.get('date', '')
    reports = data.get('reports', [])

    prs = Presentation('templateReport.pptx')

    # Slide 1: Cover
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame and "Update" in shape.text_frame.text:
            shape.text_frame.text = f"Update {report_date}"

    # Slide 2 and 3 are templates
    slide2 = prs.slides[1]
    slide4 = prs.slides[3] if len(prs.slides) > 3 else None

    # Generate slides for each user report
    for report in reports:
        new_slide = clone_slide(prs, slide2)
        
        # We need to process shapes on the cloned slide
        # Note: python-pptx will discover the cloned shapes appended to _spTree
        for shape in list(new_slide.shapes):
            if not shape.has_text_frame:
                continue
            
            text = shape.text_frame.text.strip()
            
            if "ชื่องาน" in text:
                shape.text_frame.text = report.get('title', 'Daily Report')
            elif "วันที่ดำเนินการ" in text:
                shape.text_frame.text = report.get('date', report_date)
            elif "เนื้อหา" in text:
                # Format bullet points
                bullets = report.get('summary', [])
                shape.text_frame.text = ""
                tf = shape.text_frame
                tf.word_wrap = True
                for index, bullet in enumerate(bullets):
                    p = tf.add_paragraph() if index > 0 else tf.paragraphs[0]
                    p.text = f"• {bullet}"
            elif "รูปประกอบ" in text:
                # Get the position coordinates of the placeholder shape
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Remove the placeholder shape
                new_slide.shapes._spTree.remove(shape.element)
                
                # Add the actual image if present
                img_base64 = report.get('base64Image')
                if img_base64:
                    try:
                        header, encoded = img_base64.split(",", 1) if "," in img_base64 else ("", img_base64)
                        img_data = base64.b64decode(encoded)
                        img_stream = io.BytesIO(img_data)
                        new_slide.shapes.add_picture(img_stream, left, top, width, height)
                    except Exception as img_err:
                        print(f"Error adding image for report: {img_err}")
                else:
                    # Put a text block indicating no image
                    txBox = new_slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = "[ไม่มีรูปประกอบ]"

    # Clone the closing slide (Slide 4) to the end
    if slide4:
        clone_slide(prs, slide4)

    # Delete the template slides (Slide 2 and Slide 3) and the original Slide 4
    # Initial list has 4 template slides at start: [S1, S2, S3, S4, ...new_slides...]
    # S4 is at index 3, S3 is at index 2, S2 is at index 1
    # We must delete in reverse order to keep indices correct:
    id_list = prs.slides._sldIdLst
    if len(id_list) > 3:
        del id_list[3]  # Original S4
    if len(id_list) > 2:
        del id_list[2]  # Original S3
    if len(id_list) > 1:
        del id_list[1]  # Original S2

    prs.save(output_path)
    print("PPTX generated successfully at:", output_path)

if __name__ == "__main__":
    main()
