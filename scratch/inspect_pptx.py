from pptx import Presentation

try:
    prs = Presentation('templateReport.pptx')
    print(f'Total slides: {len(prs.slides)}')
    for i, slide in enumerate(prs.slides):
        print(f'\n--- Slide {i+1} ---')
        for j, shape in enumerate(slide.shapes):
            text = ''
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
            print(f'Shape {j+1}: Name={shape.name}, Type={shape.shape_type}, Text={text[:100]}')
except Exception as e:
    print('Error:', e)
