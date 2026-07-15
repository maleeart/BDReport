import sys
from pptx import Presentation

# Write output in UTF-8
with open('scratch/inspect_utf8.txt', 'w', encoding='utf-8') as f:
    try:
        prs = Presentation('templateReport.pptx')
        f.write(f'Total slides: {len(prs.slides)}\n')
        for i, slide in enumerate(prs.slides):
            f.write(f'\n--- Slide {i+1} ---\n')
            for j, shape in enumerate(slide.shapes):
                text = ''
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                f.write(f'Shape {j+1}: Name="{shape.name}", Type={shape.shape_type}, Text="{text}"\n')
    except Exception as e:
        f.write(f'Error: {e}\n')
print('Done!')
