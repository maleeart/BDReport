from http.server import BaseHTTPRequestHandler
import json
import urllib.parse
import urllib.request
import io
import base64
import copy
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image


def pick_layout(aspects, width, height, gap):
    """Pick the row split that renders images as large as possible.

    Tries every balanced split (r rows, images spread evenly) and keeps the one
    with the biggest row height. Rows are packed by the images' real aspect
    ratios, so no image is letterboxed inside an oversized cell.
    Returns (images_per_row, row_height).
    """
    best = None
    for r in range(1, len(aspects) + 1):
        base, rem = divmod(len(aspects), r)
        counts = [base + 1] * rem + [base] * (r - rem)
        h = (height - (r - 1) * gap) / r
        i = 0
        for c in counts:
            h = min(h, (width - (c - 1) * gap) / sum(aspects[i:i + c]))
            i += c
        if h > 0 and (best is None or h > best[1]):
            best = (counts, h)
    return best

def disable_bullets(p):
    try:
        pPr = p._p.get_or_add_pPr()
        for child in list(pPr):
            if any(tag in child.tag for tag in ['buChar', 'buAutoNum', 'buSzPct', 'buFont', 'buNone']):
                pPr.remove(child)
        buNone = OxmlElement('a:buNone')
        pPr.append(buNone)
    except Exception as e:
        print(f"Error disabling bullets: {e}")

def clone_slide(prs, src_slide):
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    for shape in src_slide.shapes:
        new_shape_xml = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_shape_xml)
    
    # Copy relationships to preserve template images/graphics
    for rId, rel in src_slide.part.rels._rels.items():
        if rId not in new_slide.part.rels._rels:
            new_slide.part.rels._rels[rId] = rel
        elif rel.reltype != "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout":
            new_slide.part.rels._rels[rId] = rel
    return new_slide

class handler(BaseHTTPRequestHandler):
    def do_GET(self):
        try:
            # 1. Parse Query Parameters
            parsed_url = urllib.parse.urlparse(self.path)
            query_params = urllib.parse.parse_qs(parsed_url.query)
            
            query_secret = query_params.get('secret', [None])[0]
            week_param = query_params.get('week', [None])[0]
            indices_param = query_params.get('indices', [None])[0]
            group_id_param = query_params.get('groupId', [None])[0]
            
            # Default to current week if not provided (adjusted to Bangkok timezone UTC+7)
            if not week_param:
                from datetime import datetime, timedelta
                local_now = datetime.utcnow() + timedelta(hours=7)
                iso_year, iso_week, _ = local_now.isocalendar()
                week_param = f"{iso_year}-W{iso_week:02d}"

            # 2. Check Authentication
            cron_secret = os.environ.get('CRON_SECRET')
            auth_header = self.headers.get('Authorization')
            
            is_authorized = not cron_secret or \
                            (auth_header == f"Bearer {cron_secret}") or \
                            (query_secret == cron_secret)

            if not is_authorized:
                self.send_response(401)
                self.send_header('Content-type', 'application/json')
                self.end_headers()
                self.wfile.write(json.dumps({'error': 'Unauthorized'}).encode('utf-8'))
                return

            # 3. Fetch Report Data from local Node.js API
            host = self.headers.get('Host', 'localhost:3000')
            protocol = 'https' if 'https' in self.headers.get('X-Forwarded-Proto', '') else 'http'
            reports_url = f"{protocol}://{host}/api/reports?week={week_param}&includeImages=true"
            if group_id_param:
                reports_url += f"&groupId={urllib.parse.quote(group_id_param)}"

            # Make request
            req = urllib.request.Request(reports_url)
            try:
                with urllib.request.urlopen(req) as response:
                    res_body = response.read().decode('utf-8')
                    report_data = json.loads(res_body)
            except Exception as fetch_err:
                self.send_response(500)
                self.send_header('Content-type', 'application/json')
                self.end_headers()
                self.wfile.write(json.dumps({'error': f'Failed to fetch reports: {str(fetch_err)}'}).encode('utf-8'))
                return

            # If no reports, return JSON info
            reports = report_data.get('reports', [])
            report_date = report_data.get('date', '')

            # If indices are specified, select them from the original unfiltered reports first
            if indices_param is not None:
                try:
                    selected_indices = [int(x) for x in indices_param.split(',') if x.strip().isdigit()]
                    reports = [reports[i] for i in selected_indices if i < len(reports)]
                except Exception as filter_err:
                    print(f"Error filtering reports by indices: {filter_err}")
            # Otherwise, if a specific groupId is requested, filter reports for that group (with fallback if empty)
            elif group_id_param and group_id_param != 'all':
                matching_reports = [r for r in reports if r.get('groupId') == group_id_param]
                if matching_reports:
                    reports = matching_reports

            if not reports:
                self.send_response(200)
                self.send_header('Content-type', 'application/json')
                self.end_headers()
                self.wfile.write(json.dumps({'message': 'No reports found for this date', 'date': report_date}).encode('utf-8'))
                return

            # 4. Open template and modify PPTX
            # templateReport.pptx should be located in the root of the project
            # Vercel deploys files relative to the project root
            template_path = 'templateReport.pptx'
            if not os.path.exists(template_path):
                # Try finding in parent directory if Vercel nesting differs
                template_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'templateReport.pptx')
                if not os.path.exists(template_path):
                    template_path = 'templateReport.pptx'

            prs = Presentation(template_path)

            # Slide 1: Cover
            slide1 = prs.slides[0]
            for shape in slide1.shapes:
                if shape.has_text_frame and "Update" in shape.text_frame.text:
                    shape.text_frame.text = f"Update {report_date}"
                    p = shape.text_frame.paragraphs[0]
                    if len(p.runs) > 0:
                        p.runs[0].font.name = "TH Sarabun New"
                        p.runs[0].font.size = Pt(18)
                        p.runs[0].font.bold = True

            # Slide 2 and 3 are templates
            slide2 = prs.slides[1]
            slide4 = prs.slides[3] if len(prs.slides) > 3 else None

            # Generate slides for each user report
            for report in reports:
                new_slide = clone_slide(prs, slide2)
                
                # Delete all extra/redundant shapes from the cloned slide
                for shape in list(new_slide.shapes):
                    keep = False
                    if shape.has_text_frame:
                        t_text = shape.text_frame.text.strip()
                        if "ชื่องาน" in t_text or t_text.startswith("งาน") or "วันที่ดำเนินการ" in t_text or "เนื้อหา" in t_text or "รูปประกอบ" in t_text:
                            keep = True
                    if not keep:
                        try:
                            new_slide.shapes._spTree.remove(shape.element)
                        except Exception:
                            pass

                for shape in list(new_slide.shapes):
                    if not shape.has_text_frame:
                        continue
                    
                    text = shape.text_frame.text.strip()
                    
                    # Handle subtitle header and main task title on Slide 2
                    if "ชื่องาน" in text or text.startswith("งาน"):
                        is_main_title = "ชื่องาน" in text or text.startswith("งาน")
                        # Capture original run formatting if present
                        src_run = shape.text_frame.paragraphs[0].runs[0] if (len(shape.text_frame.paragraphs) > 0 and len(shape.text_frame.paragraphs[0].runs) > 0) else None
                        
                        shape.text_frame.text = report.get('title', 'Weekly Report')
                        p = shape.text_frame.paragraphs[0]
                        disable_bullets(p) # Completely disable bullet points on title
                        if len(p.runs) > 0:
                            run = p.runs[0]
                            run.font.name = "TH Sarabun New"
                            run.font.size = Pt(32) if is_main_title else Pt(12)
                            run.font.bold = True
                            if src_run and src_run.font.color:
                                try:
                                    if src_run.font.color.type == 1:
                                        run.font.color.rgb = src_run.font.color.rgb
                                    elif src_run.font.color.type == 2:
                                        run.font.color.theme_color = src_run.font.color.theme_color
                                except Exception:
                                    pass
                    elif "วันที่ดำเนินการ" in text:
                        shape.text_frame.text = report.get('date', report_date)
                        p = shape.text_frame.paragraphs[0]
                        disable_bullets(p) # Completely disable bullet points on date
                        p.alignment = 3 # Right align
                        if len(p.runs) > 0:
                            run = p.runs[0]
                            run.font.name = "TH Sarabun New"
                            run.font.size = Pt(14)
                    elif "เนื้อหา" in text:
                        bullets = report.get('summary', [])
                        
                        tf = shape.text_frame
                        tf.word_wrap = True
                        
                        # Preserve original template formatting if possible
                        first_p = tf.paragraphs[0]
                        src_run = first_p.runs[0] if len(first_p.runs) > 0 else None
                        
                        # Clear default placeholder text from the template's first paragraph
                        first_p.text = ""
                        
                        # Remove other paragraphs from the placeholder if any
                        while len(tf.paragraphs) > 1:
                            p_el = tf.paragraphs[1]._p
                            p_el.getparent().remove(p_el)
                            
                        # Populate bullets with native Level 0 styling
                        for index, bullet in enumerate(bullets):
                            p = tf.add_paragraph() if index > 0 else tf.paragraphs[0]
                            p.level = 0
                            p.text = bullet
                            
                            # Apply unified premium Thai font settings
                            if len(p.runs) > 0:
                                run = p.runs[0]
                                run.font.name = "TH Sarabun New"
                                run.font.size = Pt(24)
                                if src_run:
                                    run.font.bold = src_run.font.bold
                                    if src_run.font.color:
                                        try:
                                            if src_run.font.color.type == 1: # RGB
                                                run.font.color.rgb = src_run.font.color.rgb
                                            elif src_run.font.color.type == 2: # Theme Color
                                                run.font.color.theme_color = src_run.font.color.theme_color
                                        except Exception:
                                            pass
                    elif "รูปประกอบ" in text:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        new_slide.shapes._spTree.remove(shape.element)
                        
                        img_list = report.get('base64Images', [])
                        if not img_list and report.get('base64Image'):
                            img_list = [report.get('base64Image')]
                            
                        # Decode every image once, keeping its real aspect ratio
                        imgs = []
                        for idx, img_base64 in enumerate(img_list):
                            try:
                                encoded = img_base64.split(",", 1)[1] if "," in img_base64 else img_base64
                                img_stream = io.BytesIO(base64.b64decode(encoded))
                                size = Image.open(img_stream).size
                                imgs.append((img_stream, size[0] / size[1]))
                            except Exception as img_err:
                                print(f"Skipping unreadable image {idx}: {img_err}")

                        if imgs:
                            gap = 100000  # 100,000 EMUs gap

                            counts, row_h = pick_layout([a for _, a in imgs], width, height, gap)

                            rows = len(counts)
                            block_height = rows * row_h + (rows - 1) * gap
                            y_start = top + (height - block_height) / 2

                            i = 0
                            for r_idx, cols_in_row in enumerate(counts):
                                row = imgs[i:i + cols_in_row]
                                i += cols_in_row

                                row_width = sum(a * row_h for _, a in row) + (cols_in_row - 1) * gap
                                x = left + (width - row_width) / 2
                                row_top = y_start + r_idx * (row_h + gap)

                                for img_stream, aspect in row:
                                    img_w = aspect * row_h
                                    try:
                                        img_stream.seek(0)
                                        new_slide.shapes.add_picture(img_stream, int(x), int(row_top), int(img_w), int(row_h))
                                    except Exception as img_err:
                                        print(f"Error adding image: {img_err}")
                                    x += img_w + gap
                        else:
                            txBox = new_slide.shapes.add_textbox(left, top, width, height)
                            txBox.text_frame.text = "[ไม่มีรูปประกอบ]"

            # Move the original Slide 4 (Closing) to the end of the presentation
            # This preserves its original media relationships perfectly
            id_list = prs.slides._sldIdLst
            if slide4 and len(id_list) > 3:
                slide4_id = id_list[3]
                id_list.remove(slide4_id)
                id_list.append(slide4_id)

            # Delete the template slides (Slide 2 and Slide 3)
            # Reverse order (index 2 first, then index 1)
            if len(id_list) > 2:
                del id_list[2]  # Original S3
            if len(id_list) > 1:
                del id_list[1]  # Original S2

            # 5. Save PPTX into an in-memory stream
            output_stream = io.BytesIO()
            prs.save(output_stream)
            output_stream.seek(0)
            pptx_bytes = output_stream.read()

            # 6. Optional: Upload to Discord if Webhook is set
            discord_url = os.environ.get('DISCORD_WEBHOOK_URL')
            if discord_url:
                try:
                    # Construct multipart form-data payload in python without external libraries
                    boundary = '----BDReportBoundary'
                    payload = []
                    
                    # File field
                    payload.append(f'--{boundary}')
                    payload.append(f'Content-Disposition: form-data; name="file"; filename="report-{week_param}.pptx"')
                    payload.append('Content-Type: application/vnd.openxmlformats-officedocument.presentationml.presentation')
                    payload.append('')
                    payload.append(pptx_bytes)
                    
                    # Payload JSON field
                    payload.append(f'--{boundary}')
                    payload.append('Content-Disposition: form-data; name="payload_json"')
                    payload.append('Content-Type: application/json')
                    payload.append('')
                    payload.append(json.dumps({
                        'content': f'📊 **BDReport PowerPoint Generated from Template**\nDate: {report_date}'
                    }))
                    
                    payload.append(f'--{boundary}--')
                    payload.append('')
                    
                    # Merge payload parts
                    body = b''
                    for part in payload:
                        if isinstance(part, str):
                            body += part.encode('utf-8') + b'\r\n'
                        else:
                            body += part + b'\r\n'
                            
                    discord_req = urllib.request.Request(
                        discord_url,
                        data=body,
                        headers={
                            'Content-Type': f'multipart/form-data; boundary={boundary}',
                            'Content-Length': len(body)
                        },
                        method='POST'
                    )
                    with urllib.request.urlopen(discord_req) as discord_res:
                        pass
                except Exception as discord_err:
                    print(f"Failed to post to Discord: {discord_err}")

            # 7. Send Response File
            self.send_response(200)
            self.send_header('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            # Parse week number and year
            year = "2026"
            week_no = week_param
            if "-W" in week_param:
                parts = week_param.split("-W")
                year = parts[0]
                week_no = parts[1]
            
            filename = f"Weekly Report (EGAT IOT) (week {week_no}-{year}).pptx"
            safe_filename = urllib.parse.quote(filename)
            self.send_header('Content-Disposition', f"attachment; filename*=UTF-8''{safe_filename}")
            self.send_header('Content-Length', len(pptx_bytes))
            self.end_headers()
            self.wfile.write(pptx_bytes)

        except Exception as e:
            self.send_response(500)
            self.send_header('Content-type', 'application/json')
            self.end_headers()
            self.wfile.write(json.dumps({'error': str(e)}).encode('utf-8'))
            return


if __name__ == "__main__":
    # Self-check for pick_layout: a wide 16:9 image area, 100k EMU gaps.
    W, H, GAP = 9144000, 4000000, 100000

    # The regression: 9+ images must wrap into rows, not one tiny single row.
    counts, h = pick_layout([1.0] * 9, W, H, GAP)
    assert len(counts) > 1 and h > (W - 8 * GAP) / 9, (counts, h)

    # More images must never render them larger.
    prev = None
    for n in range(1, 21):
        _, h = pick_layout([1.0] * n, W, H, GAP)
        assert h > 0
        assert prev is None or h <= prev, n
        prev = h

    # Rows are packed by real aspect ratios: a row must fit the area exactly.
    aspects = [1.78, 0.75, 1.0, 1.33, 1.78, 0.56, 1.0, 1.33, 1.78, 1.0]
    counts, h = pick_layout(aspects, W, H, GAP)
    assert sum(counts) == len(aspects)
    assert len(counts) * h + (len(counts) - 1) * GAP <= H + 1
    i = 0
    for c in counts:
        assert sum(a * h for a in aspects[i:i + c]) + (c - 1) * GAP <= W + 1
        i += c

    print("pick_layout self-check passed")
